# -*- coding: utf-8 -*-
"""
Build the PhysicsMind EXIST-consultation deck as an editable .pptx.
Mirrors EXIST_deck.html (8 slides, B2C, dark product theme) and embeds the
real Lorentz-force product screenshot on the Solution slide.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# ---- palette (from EXIST_deck.html :root) -------------------------------
BG    = RGBColor(0x0F, 0x0B, 0x08)   # dark canvas
INK   = RGBColor(0xF4, 0xEC, 0xE1)
MUTED = RGBColor(0xBC, 0xAB, 0x95)
FAINT = RGBColor(0x8C, 0x7C, 0x69)
ACC   = RGBColor(0xE0, 0x7A, 0x42)   # terracotta-orange
ACC2  = RGBColor(0xEA, 0xA6, 0x3C)   # amber/gold
BLUE  = RGBColor(0x5B, 0x8B, 0xC4)
GREEN = RGBColor(0x6F, 0xCF, 0x7F)
BODY  = RGBColor(0xEC, 0xE3, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0x1C, 0x14, 0x0D)

SERIF = "Georgia"
SANS  = "Segoe UI"

IMG = r"C:\Users\PRADEEEP\.claude\image-cache\507b8370-f63a-4e68-8de3-aec2f6c93ee9\1.png"
OUT = r"C:\Tutor\physics-mind\docs\pitch\PhysicsMind_EXIST_2026-06-22_v3.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
TOTAL = 9

# ---- helpers ------------------------------------------------------------
def _noshadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass

def add_text(slide, l, t, w, h, paragraphs, align=PP_ALIGN.LEFT,
             line_spacing=1.0, space_after=0, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paragraphs = list of paragraphs; each paragraph = list of runs;
       each run = (text, bold, color, font, size)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        for (text, bold, color, font, size) in runs:
            r = p.add_run()
            r.text = text
            r.font.bold = bold
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = font
    return tb

def set_tracking(tb, spc=260):
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            r._r.get_or_add_rPr().set('spc', str(spc))

def set_italic(tb):
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            r.font.italic = True

def kicker(slide, text, l=Inches(0.9), t=Inches(1.18)):
    tb = add_text(slide, l, t, Inches(9), Inches(0.4),
                  [[(text.upper(), True, ACC, SANS, 12)]])
    set_tracking(tb, 280)
    return tb

def headline(slide, text, l=Inches(0.9), t=Inches(1.66), w=Inches(11.4),
             h=Inches(1.7), size=31, color=INK):
    return add_text(slide, l, t, w, h, [[(text, True, color, SERIF, size)]],
                    line_spacing=1.08)

def bullets(slide, l, t, w, h, items, marker=ACC, size=16.5, gap=11):
    paras = []
    for segs in items:
        runs = [("◆   ", False, marker, SANS, size - 3)]
        for (txt, bold) in segs:
            runs.append((txt, bold, WHITE if bold else BODY, SANS, size))
        paras.append(runs)
    return add_text(slide, l, t, w, h, paras, line_spacing=1.12, space_after=gap)

def caption(slide, text, t=Inches(5.75), l=Inches(0.92), w=Inches(10), color=FAINT):
    tb = add_text(slide, l, t, w, Inches(0.7), [[(text, False, color, SANS, 13)]])
    set_italic(tb)
    return tb

def pill(slide, l, t, w, h, text, color=ACC, size=11):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.background()
    sh.line.color.rgb = color
    sh.line.width = Pt(1.25)
    _noshadow(sh)
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.bold = True; r.font.size = Pt(size); r.font.name = SANS; r.font.color.rgb = color
    set_tracking(sh, 180)
    return sh

def chrome(slide, idx):
    # progress bar
    bw = int(SW * (idx / TOTAL))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, bw, Pt(3.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACC; bar.line.fill.background()
    _noshadow(bar)
    # brand bottom-left
    add_text(slide, Inches(0.5), Inches(6.98), Inches(4), Inches(0.4),
             [[("Physics", True, ACC, SERIF, 12.5), ("Mind", True, INK, SERIF, 12.5)]])
    # counter bottom-right
    add_text(slide, SW - Inches(1.9), Inches(6.98), Inches(1.4), Inches(0.4),
             [[("%d / %d" % (idx, TOTAL), False, FAINT, SANS, 11)]], align=PP_ALIGN.RIGHT)

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def new_slide(idx):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    chrome(s, idx)
    return s

# ---- SLIDE 1 — Title ----------------------------------------------------
s = new_slide(1)
kicker(s, "Deep-tech education · Dresden", t=Inches(1.35))
# logomark
logo = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.12),
                          Inches(0.98), Inches(0.98))
logo.fill.solid(); logo.fill.fore_color.rgb = ACC; logo.line.fill.background()
_noshadow(logo)
# atom: three ellipses + nucleus
cx, cy = Inches(0.9) + Inches(0.49), Inches(2.12) + Inches(0.49)
ow, oh = Inches(0.82), Inches(0.30)
for ang in (0, 60, 120):
    el = s.shapes.add_shape(MSO_SHAPE.OVAL, int(cx - ow/2), int(cy - oh/2), ow, oh)
    el.fill.background(); el.line.color.rgb = WHITE; el.line.width = Pt(1.4)
    el.rotation = ang; _noshadow(el)
nuc = s.shapes.add_shape(MSO_SHAPE.OVAL, int(cx - Inches(0.075)), int(cy - Inches(0.075)),
                         Inches(0.15), Inches(0.15))
nuc.fill.solid(); nuc.fill.fore_color.rgb = WHITE; nuc.line.fill.background(); _noshadow(nuc)
# wordmark
add_text(s, Inches(2.08), Inches(2.02), Inches(9), Inches(1.18),
         [[("Physics", True, INK, SERIF, 52), ("Mind", True, ACC, SERIF, 52)]],
         anchor=MSO_ANCHOR.MIDDLE)
# tagline
tag = add_text(s, Inches(0.92), Inches(3.42), Inches(10.5), Inches(0.8),
               [[("A personal science tutor for every student who doesn’t have one.",
                  False, MUTED, SERIF, 22)]])
set_italic(tag)
# problem → promise contrast
prob = add_text(s, Inches(0.95), Inches(4.14), Inches(11.4), Inches(0.42),
                [[("A great teacher can’t be there 24×7 — so learning stays passive.",
                   False, MUTED, SANS, 15)]])
set_italic(prob)
add_text(s, Inches(0.95), Inches(4.56), Inches(11.4), Inches(0.42),
         [[("PhysicsMind is ", False, INK, SANS, 15),
           ("always on", True, ACC, SANS, 15),
           (" — learning that’s ", False, INK, SANS, 15),
           ("active and proactive", True, ACC, SANS, 15),
           (", never passive.", False, INK, SANS, 15)]])
# meta
add_text(s, Inches(0.95), Inches(5.24), Inches(11), Inches(1.0),
         [[("Pradeep Nagapuri", True, INK, SANS, 15),
           ("  —  Graduate, Dresden International University (2026)", False, MUTED, SANS, 15)],
          [("Deep-tech education startup  ·  built in Dresden, Germany", False, MUTED, SANS, 15)]],
         line_spacing=1.4, space_after=4)
# badge
pill(s, Inches(0.9), Inches(6.4), Inches(5.05), Inches(0.5),
     "APPLYING FOR EXIST-GRÜNDERSTIPENDIUM", color=ACC, size=11)
notes(s, "Say: “I just finished my degree at Dresden International University. I'm building "
         "a deep-tech education company — an always-on AI tutor that makes students learn "
         "actively and proactively, not just watch passively — and I'd like to understand "
         "whether and how EXIST can support it.”")

# ---- SLIDE 2 — Problem --------------------------------------------------
s = new_slide(2)
kicker(s, "The Problem")
headline(s, "Understanding science shouldn’t depend on affording a great teacher.",
         w=Inches(10.5), size=30)
bullets(s, Inches(0.9), Inches(3.5), Inches(11), Inches(3),
        [[("A great tutor is the biggest advantage in learning — and ", False),
          ("never available 24×7", True)],
         [("Video and text are ", False), ("passive", True),
          (" — students watch, they don’t do", False)],
         [("Students memorise formulas; ", False), ("misconceptions go uncorrected", True)],
         [("1-to-1 tutoring is expensive and ", False), ("doesn’t scale", True)]])
notes(s, "Say: “The real gap isn't content — it's access to a teacher who makes it click "
         "and corrects how you think.”")

# ---- SLIDE 3 — Solution (with product screenshot) -----------------------
s = new_slide(3)
kicker(s, "The Solution")
headline(s, "We make the concept move.", w=Inches(6.2), size=34)
bullets(s, Inches(0.9), Inches(3.0), Inches(6.0), Inches(3.6),
        [[("Physics-accurate, ", False), ("interactive simulations", True),
          (" — learn by doing", False)],
         [("An ", False), ("AI professor", True),
          (" that explains & drives the sim live, by voice", False)],
         [("Always on — ", True), ("24×7 doubt-solving", True),
          (", no waiting for a teacher", False)],
         [("“ChatGPT tells you — PhysicsMind shows you”", True)]],
        size=15.5)
# product screenshot on the right, with accent border
iw, ih = Image.open(IMG).size
maxw, maxh = Inches(5.55), Inches(4.35)
w = maxw
h = int(w * ih / iw)
if h > maxh:
    h = maxh
    w = int(h * iw / ih)
left = SW - w - Inches(0.85)
top = Inches(1.75)
pic = s.shapes.add_picture(IMG, left, top, width=w, height=h)
pic.line.color.rgb = ACC
pic.line.width = Pt(1.5)
add_text(s, left, top + h + Inches(0.08), w, Inches(0.4),
         [[("Live product — Lorentz force, with the voice AI professor", False, FAINT, SANS, 11)]],
         align=PP_ALIGN.CENTER)
notes(s, "Say: open the real simulation on your laptop. The moving, correct physics does more "
         "than any slide — and it's always on, so the student learns actively the moment they're "
         "stuck. Tap the mic and ask a 'what if' to show the AI professor driving it.")

# ---- SLIDE 4 — Why it's different ---------------------------------------
s = new_slide(4)
kicker(s, "Why It’s Different")
headline(s, "Not another tutoring app — a deep-tech engine.", w=Inches(10.5), size=30)
bullets(s, Inches(0.9), Inches(3.4), Inches(11.2), Inches(3),
        [[("Physics engines compute everything; the AI writes only ", False),
          ("configuration, never the physics", True), (" → always accurate", False)],
         [("A ", False), ("curriculum-specific simulation library", True), (" — our core IP", False)],
         [("A proprietary ", False), ("data flywheel", True),
          (": every student doubt improves the next explanation", False)]])
caption(s, "Defensible technology + knowledge — not a feature a competitor can copy.")
notes(s, "Say: “The innovation is the architecture and the data — that's what makes it "
         "defensible, not a feature anyone can clone.” This is the EXIST-critical slide.")

# ---- SLIDE 5 — Product & business model (B2C) ---------------------------
s = new_slide(5)
kicker(s, "Product · B2C")
headline(s, "Direct-to-student subscription.", size=32)
bullets(s, Inches(0.9), Inches(3.4), Inches(11), Inches(3),
        [[("Students ", False), ("subscribe directly", True),
          (" — web app + AI voice professor", False)],
         [("Tiered: simulations → AI doubt-solving → voice professor", False)],
         [("Concept clarity first", True), ("; exam depth grows with the student", False)]])
caption(s, "No dependence on schools or institutions.")
notes(s, "Say: “We sell directly to students. We're not dependent on schools or "
         "institutions.”")

# ---- SLIDE 6 — Market & GTM ---------------------------------------------
s = new_slide(6)
kicker(s, "Market & Go-to-Market")
headline(s, "A focused entry into a global market.", w=Inches(10.5), size=30)
bullets(s, Inches(0.9), Inches(3.35), Inches(11.4), Inches(3),
        [[("Global STEM learning — underserved by ", False), ("simulation-first", True),
          (" products", False)],
         [("Build on the ", False), ("Indian curriculum", True),
          (" (deepest → ports to AP / A-level / IB)", False)],
         [("First revenue: ", False), ("UAE diaspora", True),
          (" — same curriculum, high willingness to pay", False)],
         [("Year 2: India scale + Western curricula", False)]])
notes(s, "Say: “One curriculum, sequenced markets. We build where our expertise is, "
         "monetize first where families pay, expand deliberately.”")

# ---- SLIDE 7 — Traction & feasibility -----------------------------------
s = new_slide(7)
kicker(s, "Traction & Feasibility")
headline(s, "Already building. Already validated by a teacher.", w=Inches(10.5), size=30)
bullets(s, Inches(0.9), Inches(3.3), Inches(11.4), Inches(2.6),
        [[("Gold-standard simulations built", True), (" — magnetism & electrostatics", False)],
         [("Working ", False), ("AI authoring pipeline", True), (" + AI voice-professor prototype", False)],
         [("A ", False), ("teacher co-development partner", True), (" reviewing real lessons", False)],
         [("Built solo, from Germany, with AI-assisted development", False)]])
pill(s, Inches(0.9), Inches(5.5), Inches(2.55), Inches(0.46), "✓  TEACHER-VERIFIED",
     color=ACC, size=10.5)
caption(s, "Pre-revenue today — a teacher-supervised student pilot is the next step.",
        t=Inches(6.12))
notes(s, "Say honestly: “Pre-revenue today — the next step is a teacher-supervised "
         "student pilot. That's exactly what the EXIST year would fund.”")

# ---- SLIDE 8 — Vision ---------------------------------------------------
s = new_slide(8)
kicker(s, "The Vision")
headline(s, "A great teacher for every student — always on, in every subject.",
         w=Inches(11.4), size=30)
bullets(s, Inches(0.9), Inches(3.35), Inches(11.4), Inches(3),
        [[("Every concept in science", True),
          (" becomes an interactive, physics-accurate simulation", False)],
         [("An ", False), ("AI professor", True),
          (" that teaches, adapts, and answers by voice — ", False), ("24×7", True)],
         [("Physics first → ", False),
          ("chemistry, biology, mathematics", True), (" next", False)],
         [("Indian curriculum → ", False), ("AP · A-level · IB", True),
          (" → students everywhere", False)]])
pill(s, Inches(0.9), Inches(5.55), Inches(4.7), Inches(0.5),
     "ALWAYS ON · ACTIVE · PROACTIVE", color=ACC, size=11)
caption(s, "From “who can afford a great teacher” to “every student has one.”",
        t=Inches(6.2))
notes(s, "Say: “This is where it goes. The same engine that teaches Lorentz force today "
         "teaches every concept in science tomorrow — one AI professor, always on, for "
         "every student, in every curriculum. That's the company I want to build with EXIST.”")

# ---- SLIDE 9 — EXIST year + the ask -------------------------------------
s = new_slide(9)
kicker(s, "EXIST · The Ask")
headline(s, "What the stipend makes possible.", size=32)
# col 1
add_text(s, Inches(0.9), Inches(3.15), Inches(5.4), Inches(0.4),
         [[("THE 12-MONTH PLAN", True, ACC2, SANS, 12)]])
set_tracking(add_text(s, Inches(0.9), Inches(3.15), Inches(5.4), Inches(0.4),
             [[("THE 12-MONTH PLAN", True, ACC2, SANS, 12)]]), 200)
bullets(s, Inches(0.9), Inches(3.62), Inches(5.4), Inches(3),
        [[("Build the concept library to depth", False)],
         [("Run a teacher + student pilot", False)],
         [("Reach first paying users", False)],
         [("Build a German-based deep-tech company", False)]],
        marker=ACC2, size=15, gap=9)
# col 2
set_tracking(add_text(s, Inches(6.85), Inches(3.15), Inches(6.0), Inches(0.4),
             [[("MY QUESTIONS FOR YOU", True, ACC2, SANS, 12)]]), 200)
bullets(s, Inches(6.85), Inches(3.62), Inches(6.0), Inches(3),
        [[("Eligibility (DIU graduate, 2026)", False)],
         [("Team requirement — do I need co-founders?", False)],
         [("Professor / mentor — can you help me find one?", False)],
         [("Visa pathway (§20 → §21)", False)],
         [("Next steps & timeline", False)]],
        marker=ACC2, size=15, gap=9)
notes(s, "Say: “I'd love your guidance on these — that's really why I'm here.” "
         "Let the conversation open up after this slide.")

try:
    prs.save(OUT)
    saved = OUT
except PermissionError:
    saved = OUT.replace(".pptx", "_v2.pptx")
    prs.save(saved)
    print("NOTE: primary file was locked (open in PowerPoint); saved to v2 instead.")
print("OK saved:", saved)
print("slides:", len(prs.slides._sldIdLst))
print("image:", "%dx%d" % (iw, ih), "-> placed", round(w/914400, 2), "x", round(h/914400, 2), "in")
