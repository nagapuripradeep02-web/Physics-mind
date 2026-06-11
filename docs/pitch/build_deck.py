"""
PhysicsMind — pitch deck generator (v2, story-driven, 10 slides).

Produces a native, fully-editable PowerPoint (.pptx) at:
    C:/Tutor/physics-mind/docs/pitch/PhysicsMind_Pitch.pptx

Re-run anytime to regenerate after editing slide content below.
    python build_deck.py

Story arc: Dream → Problem → Gap → Solution → Moat → Revenue →
           Production pace → Why a team → 3-year roadmap → Next step.
Palette aligned to the product's real identity (dark navy + field-blue +
physics accents orange/green/amber). Speaker notes on every slide.
Honest: no fake traction; the confusion-data flywheel is framed as
activating WITH the pilot, never as existing data.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------- palette (product-aligned)
BG_DARK   = RGBColor(0x0A, 0x0F, 0x1E)   # near-black navy (matches sim background)
BLUE      = RGBColor(0x1E, 0x88, 0xE5)   # PRIMARY accent (prints well)
BLUE_LT   = RGBColor(0x42, 0xA5, 0xF5)   # product field-blue (for dark slides)
BLUE_TINT = RGBColor(0xE7, 0xF1, 0xFB)   # light card fill
ORANGE    = RGBColor(0xF5, 0x7C, 0x00)   # particle / velocity accent
GREEN     = RGBColor(0x1B, 0x9E, 0x5A)   # force accent / "done"
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)   # next-step highlight
INK       = RGBColor(0x1E, 0x29, 0x3B)   # body text on light
MUTED     = RGBColor(0x64, 0x74, 0x8B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xF8, 0xFA, 0xFC)
BORDER    = RGBColor(0xE2, 0xE8, 0xF0)
ONDARK    = RGBColor(0xBB, 0xDB, 0xFB)   # light-blue text on dark
ONDARK_MU = RGBColor(0x93, 0xA4, 0xBC)   # muted text on dark
FONT      = "Segoe UI"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
ML, MR = Inches(0.9), Inches(0.9)
CW = SLIDE_W - ML - MR

prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
BLANK = prs.slide_layouts[6]

# ----------------------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color):
    f = s.background.fill; f.solid(); f.fore_color.rgb = color

def rect(s, x, y, w, h, color, radius=None):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    if radius is not None:
        try: sp.adjustments[0] = radius
        except Exception: pass
    return sp

def outline(sp, color, w=Pt(1)):
    sp.line.color.rgb = color; sp.line.width = w

def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf

def para(tf, size, color, bold=False, align=PP_ALIGN.LEFT, before=0, after=6, leading=1.05, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_before = Pt(before); p.space_after = Pt(after); p.line_spacing = leading
    return p

def run(p, text, size, color, bold=False, italic=False, font=FONT):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = color
    return r

def footer(s):
    tf = tb(s, ML, Inches(7.02), Inches(7), Inches(0.35))
    p = para(tf, 10, MUTED, first=True, after=0)
    run(p, "PhysicsMind", 10, MUTED, bold=True)
    run(p, "   ·   A physics tutor that makes the concept click", 10, MUTED)

def pagenum(s, n):
    tf = tb(s, SLIDE_W - Inches(1.4), Inches(7.02), Inches(0.5), Inches(0.35))
    p = para(tf, 10, MUTED, align=PP_ALIGN.RIGHT, first=True, after=0)
    run(p, str(n), 10, MUTED, bold=True)

def header(s, kicker, title, title_size=32, lines=1):
    rect(s, 0, 0, SLIDE_W, Inches(0.16), BLUE)
    tf = tb(s, ML, Inches(0.52), CW, Inches(0.4))
    p = para(tf, 13, BLUE, bold=True, first=True, after=0)
    run(p, kicker.upper(), 13, BLUE, bold=True)
    tf2 = tb(s, ML, Inches(0.92), CW, Inches(1.2))
    p2 = para(tf2, title_size, INK, bold=True, first=True, leading=1.0)
    run(p2, title, title_size, INK, bold=True)
    div_y = Inches(1.92) if lines == 1 else Inches(2.22)
    rect(s, ML, div_y, CW, Pt(2), BORDER)
    footer(s)
    return div_y + Inches(0.33)

def bullet(tf, text, lead_color=BLUE, size=15, after=10, sub=False, first=False):
    p = para(tf, size, INK, first=first, after=after, leading=1.08)
    if sub:
        run(p, "–  ", size, MUTED, bold=True); run(p, text, size, MUTED)
    else:
        run(p, "▸  ", size, lead_color, bold=True); run(p, text, size, INK)
    return p

def bullet_rich(tf, segments, size=15, after=10, first=False, lead_color=BLUE):
    p = para(tf, size, INK, first=first, after=after, leading=1.08)
    run(p, "▸  ", size, lead_color, bold=True)
    for text, color, bold in segments:
        run(p, text, size, color, bold=bold)
    return p

def card(s, x, y, w, h, fill=LIGHT, border=BORDER, radius=0.06):
    sp = rect(s, x, y, w, h, fill, radius=radius)
    if border is not None: outline(sp, border, Pt(1))
    return sp

def stat_card(s, x, y, w, big, label, accent=BLUE, h=Inches(1.5), big_size=30):
    card(s, x, y, w, h, fill=LIGHT, border=BORDER)
    tf = tb(s, x + Inches(0.22), y + Inches(0.18), w - Inches(0.44), h - Inches(0.36), anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, big_size, accent, bold=True, first=True, after=4, leading=0.95)
    run(p, big, big_size, accent, bold=True)
    p2 = para(tf, 12.5, MUTED, after=0, leading=1.05)
    run(p2, label, 12.5, MUTED)

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

# ============================================================================= SLIDE 1 — TITLE
s = slide(); bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.22), SLIDE_H, BLUE)
tf = tb(s, Inches(1.1), Inches(1.85), Inches(11), Inches(1.4))
p = para(tf, 58, WHITE, bold=True, first=True, after=0, leading=0.98)
run(p, "PhysicsMind", 58, WHITE, bold=True)
tf2 = tb(s, Inches(1.12), Inches(3.05), Inches(10.8), Inches(1.6))
p = para(tf2, 23, ONDARK, first=True, after=6, leading=1.12)
run(p, "A personal physics tutor for every student who doesn't have one.", 23, BLUE_LT)
p = para(tf2, 15, ONDARK_MU, after=0, leading=1.25)
run(p, "Interactive simulations that make a hard concept click in one viewing — for India's Class 10–12, JEE / NEET / board students.",
    15, ONDARK_MU)
rect(s, Inches(1.12), Inches(6.0), Inches(3.2), Pt(2), BLUE)
tf3 = tb(s, Inches(1.12), Inches(6.2), Inches(11), Inches(0.8))
p = para(tf3, 13, RGBColor(0xCB,0xD5,0xE1), first=True, after=0)
run(p, "Pradeep  ·  Founder", 13, RGBColor(0xCB,0xD5,0xE1), bold=True)
run(p, "      Pitch overview · June 2026", 13, ONDARK_MU)
notes(s, "STORY OPEN (30 sec): Every Indian student dreams of cracking JEE or NEET, but most can't afford a great "
         "personal teacher. I'm building the next best thing — a physics tutor that makes each concept genuinely "
         "click. Let me walk you through it as a story: the problem, what we built, how it makes money, how we build "
         "it, and where it goes over three years. I'll show you a live simulation in a minute.")

# ============================================================================= SLIDE 2 — PROBLEM
s = slide(); bg(s, WHITE)
y = header(s, "01 · The problem", "Millions compete. Almost none truly understand the physics.", lines=2)
gap = Inches(0.3); cw3 = (CW - 2*gap) / 3
stat_card(s, ML,             y, cw3, "~38 lakh", "students sit JEE + NEET every year", big_size=28)
stat_card(s, ML+cw3+gap,     y, cw3, "a few thousand", "IIT + top medical seats they fight for", accent=ORANGE, big_size=22)
stat_card(s, ML+2*(cw3+gap), y, cw3, "#1", "physics — the most-feared subject", accent=GREEN, big_size=28)
tf = tb(s, ML, y + Inches(1.75), CW, Inches(2.6))
bullet(tf, "Students memorise formulas they have never actually visualised — so it collapses under exam pressure.", first=True)
bullet(tf, "The shortage isn't content — it's genuine understanding. Notes, videos and tests are everywhere; the “aha” isn't.")
bullet(tf, "A great personal teacher fixes this — but almost no student can afford one.")
pagenum(s, 2)
notes(s, "Set the stakes. ~38 lakh aspirants a year for a few thousand seats. Physics is where they break — and they "
         "break because they MEMORISE instead of UNDERSTAND. The market is drowning in content; what's missing is the "
         "moment a concept clicks. That's normally a great teacher's job — and that's exactly what almost nobody can afford.")

# ============================================================================= SLIDE 3 — THE GAP
s = slide(); bg(s, WHITE)
y = header(s, "02 · The gap", "A huge market — and the one thing nobody has built.", lines=2)
gap = Inches(0.3); cw3 = (CW - 2*gap) / 3
stat_card(s, ML,             y, cw3, "<2%", "of JEE/NEET prep is online — barely touched", accent=ORANGE, big_size=30)
stat_card(s, ML+cw3+gap,     y, cw3, "$7.2B → $17.8B", "India coaching market, 2025 → 2034", big_size=24)
stat_card(s, ML+2*(cw3+gap), y, cw3, "4.46M", "paid users at PhysicsWallah — the market pays", accent=GREEN, big_size=28)
tf = tb(s, ML, y + Inches(1.75), CW, Inches(2.5))
bullet(tf, "Recorded video is passive — you watch, you don't manipulate, you don't truly grasp it.", first=True, after=9)
bullet(tf, "Open sandboxes (e.g. PhET) aren't exam-aligned, aren't in Indian context, and give no guided path to mastery.", after=9)
bullet_rich(tf, [("PhysicsWallah proved students pay — but it's ", INK, False),
                 ("B2C; it competes with the institutes, it doesn't arm them.", INK, True)], after=9)
bullet_rich(tf, [("The open lane: ", BLUE, True),
                 ("a guided, exam-deep, Indian-context simulation layer the institutes can give their students.", INK, False)], after=0)
tf2 = tb(s, ML, Inches(6.62), CW, Inches(0.35))
p = para(tf2, 10, MUTED, first=True, after=0)
run(p, "Sources: IMARC Group; Redseer / PhysicsWallah RHP (2025).", 10, MUTED, italic=True)
pagenum(s, 3)
notes(s, "The opportunity. The market is huge and growing, but online barely scratches JEE/NEET (<2%). PhysicsWallah "
         "proves students will pay — but PW is B2C, a RIVAL to coaching institutes, not a supplier. Nobody is arming "
         "the institute with a guided, exam-deep, Indian-context simulation layer. That empty lane is ours.")

# ============================================================================= SLIDE 4 — SOLUTION
s = slide(); bg(s, WHITE)
y = header(s, "03 · The solution", "One concept. One simulation. Mastery in one viewing.")
tf = tb(s, ML, y, Inches(7.0), Inches(4.3))
bullet(tf, "A growing library of “diamond” simulations — each teaches a single concept from zero to mastery.", first=True, after=12)
bullet_rich(tf, [("“ChatGPT can ", INK, False), ("tell", INK, True), (" you. PhysicsMind ", INK, False),
                 ("shows", BLUE, True), (" you.”", INK, False),
                 ("  The student changes things and watches the physics respond.", INK, False)], after=12)
bullet(tf, "Engineered so students understand on the first viewing — interactive, not a video.", after=12)
bullet_rich(tf, [("One mode, done excellently: ", INK, False),
                 ("conceptual mastery, authored to JEE/NEET depth.", BLUE, True),
                 ("  (Exam shortcuts come later, for overseas.)", MUTED, False)], after=0)
px = ML + Inches(7.4)
card(s, px, y, SLIDE_W - MR - px, Inches(3.95), fill=BG_DARK, border=None, radius=0.05)
tfp = tb(s, px + Inches(0.35), y + Inches(0.35), SLIDE_W - MR - px - Inches(0.7), Inches(3.3))
p = para(tfp, 13, GREEN, bold=True, first=True, after=8)
run(p, "● LIVE TODAY", 13, RGBColor(0x4A,0xDE,0x80), bold=True)
p = para(tfp, 19, WHITE, bold=True, after=8, leading=1.05)
run(p, "Magnetic force on a moving charge", 19, WHITE, bold=True)
p = para(tfp, 14, ONDARK, after=10, leading=1.25)
run(p, "A full 3D interactive simulation — 7 guided steps, the student drives it.", 14, ONDARK)
p = para(tfp, 13, ONDARK_MU, after=0, leading=1.2)
run(p, "Not a mock-up. I'll demo it live right now.", 13, ONDARK_MU, italic=True)
pagenum(s, 4)
notes(s, "The product. Each simulation teaches ONE concept to mastery, interactively — the line to land is “ChatGPT "
         "tells you; PhysicsMind shows you.” We deliberately do ONE thing excellently: conceptual mastery at JEE/NEET "
         "depth (we dropped board/competitive feature-modes to go faster and deeper). Then DEMO the magnetic-force sim "
         "live — that turns the claim into proof.")

# ============================================================================= SLIDE 5 — MOAT
s = slide(); bg(s, WHITE)
y = header(s, "04 · Why we win", "Cheap to serve. Built to compound.")
tf = tb(s, ML, y, CW, Inches(4.3))
bullet_rich(tf, [("Near-zero cost to serve. ", BLUE, True),
                 ("The engine is fixed; the content does the teaching. No AI runs at view-time — so one more student costs almost nothing.", INK, False)], first=True, after=12)
bullet_rich(tf, [("Every simulation compounds. ", BLUE, True),
                 ("A growing, exam-aligned, Indian-context library a competitor would have to rebuild from zero.", INK, False)], after=12)
bullet_rich(tf, [("A data flywheel that turns on with the pilot. ", BLUE, True),
                 ("Once students use it, we capture exactly where each one gets confused and improve the content every week — ground truth no one else has.", INK, False)], after=12)
bullet_rich(tf, [("Made for India. ", BLUE, True),
                 ("Anchors students recognise (a CRT screen, the aurora, BARC's cyclotron) — relatability no Western tool offers.", INK, False)], after=0)
pagenum(s, 5)
notes(s, "Why it's defensible. (1) Architecture makes serving nearly free — software margins. (2) The library "
         "compounds — each sim is permanent inventory. (3) The flywheel: once REAL students use it (starting with the "
         "pilot — be honest we don't have that data yet), we learn exactly where they get stuck and improve weekly. "
         "(4) Indian context no Western tool has. Be clear the flywheel ACTIVATES with the pilot — it's not data we already hold.")

# ============================================================================= SLIDE 6 — REVENUE MODEL
s = slide(); bg(s, WHITE)
y = header(s, "05 · Revenue model", "Sold to institutes. Margins that scale.")
tf = tb(s, ML, y, Inches(7.1), Inches(4.3))
bullet_rich(tf, [("Subscription, per academic year — ", INK, False),
                 ("sold B2B / B2B2C to institutes & tutors.", INK, True)], first=True, after=10)
bullet(tf, "Pricing per student / year:", after=4)
bullet(tf, "Board institutes:  ₹100–200", sub=True, after=4)
bullet(tf, "Integrated JEE/NEET institutes:  ₹200–400", sub=True, after=10)
bullet_rich(tf, [("White-label: ", INK, True),
                 ("co-brand “powered by PhysicsMind” by default; ", INK, False),
                 ("+20–30% to put only their name on it.", BLUE, True)], after=10)
bullet_rich(tf, [("The market: ", INK, True), ("~38 lakh aspirants a year, online <2% — wide open.", INK, False)], after=0)
px = ML + Inches(7.5)
card(s, px, y, SLIDE_W - MR - px, Inches(2.3), fill=BLUE_TINT, border=BLUE, radius=0.05)
tfm = tb(s, px + Inches(0.35), y + Inches(0.32), SLIDE_W - MR - px - Inches(0.7), Inches(1.7), anchor=MSO_ANCHOR.MIDDLE)
p = para(tfm, 50, BLUE, bold=True, first=True, after=2, leading=0.9)
run(p, "85–90%", 50, BLUE, bold=True)
p = para(tfm, 15, INK, after=0)
run(p, "gross margin — software economics", 15, INK, bold=True)
tfx = tb(s, px, y + Inches(2.55), SLIDE_W - MR - px, Inches(1.3))
p = para(tfx, 13, MUTED, first=True, after=0, leading=1.3)
run(p, "Example: 10 integrated institutes × 2,000 students × ₹300/yr  =  ", 13, MUTED)
run(p, "₹60 lakh / year", 13, INK, bold=True)
run(p, "  at near-zero serving cost.", 13, MUTED)
pagenum(s, 6)
notes(s, "How it makes money. Subscription to institutes, not a B2C price war. Don't over-claim exact pricing — say "
         "we validate it in the pilot. Headline = 85–90% margin: once a sim exists, serving it is nearly free, so each "
         "new institute is almost pure margin. The white-label upsell: their brand on it costs them extra.")

# ============================================================================= SLIDE 7 — PRODUCTION PACE
s = slide(); bg(s, WHITE)
y = header(s, "06 · How we build", "Honest about our one bottleneck: pace.")
gap = Inches(0.3); cw3 = (CW - 2*gap) / 3
stat_card(s, ML,             y, cw3, "~10 days", "to build one quality simulation — today", accent=ORANGE, big_size=28)
stat_card(s, ML+cw3+gap,     y, cw3, "~3–5 days", "once the reusable parts exist (it speeds up)", accent=BLUE, big_size=26)
stat_card(s, ML+2*(cw3+gap), y, cw3, "forever", "it keeps adapting to students (the moat)", accent=GREEN, big_size=28)
tf = tb(s, ML, y + Inches(1.75), CW, Inches(2.5))
bullet(tf, "Each simulation: design → build → put it in front of students → adapt to their feedback. ~10 days at first.", first=True)
bullet(tf, "It gets faster — the first sims in a chapter build reusable parts, so the rest drop to a few days.")
bullet_rich(tf, [("Then it never stops improving: ", INK, False),
                 ("90% stays, small tweaks adapt it to students for years. That continuous adaptation is the moat — cheap to maintain, impossible to copy.", INK, False)])
bullet_rich(tf, [("The honest truth: ", INK, True),
                 ("our bottleneck is content velocity — how fast we author quality sims — not market demand.", INK, False)], after=0)
pagenum(s, 7)
notes(s, "The honesty slide — this is what builds an investor's trust. Be straight: a great sim takes ~10 days now; it "
         "gets faster as we build reusable parts; then it improves forever from student feedback (the moat). And our "
         "real constraint is how fast we can MAKE quality sims, not whether people want them. Most founders hide their "
         "bottleneck — we're showing ours, with a plan to beat it (next slide).")

# ============================================================================= SLIDE 8 — WHY A TEAM
s = slide(); bg(s, WHITE)
y = header(s, "07 · Why a team", "Two of us prove it. A team scales it.")
tf = tb(s, ML, y, Inches(7.1), Inches(4.3))
bullet_rich(tf, [("Founder + an IIT professor ≈ ~2 simulations a month. ", INK, True),
                 ("Enough to PROVE the method — not to build a whole catalogue.", INK, False)], first=True, after=12)
bullet_rich(tf, [("Solo, full physics would take years. ", ORANGE, True),
                 ("That is exactly why we build a team — but only after the pilot proves it works.", INK, False)], after=12)
bullet_rich(tf, [("Each new pipeline (1 subject expert + 1 builder) ≈ +2–3 sims/month. ", INK, True),
                 ("Throughput compounds; quality rises week by week.", INK, False)], after=12)
bullet_rich(tf, [("A Year-2 team (~6 people) ≈ ₹50–80 lakh/year ", INK, True),
                 ("— capital-efficient by global standards.", INK, False)], after=0)
px = ML + Inches(7.5)
card(s, px, y, SLIDE_W - MR - px, Inches(2.7), fill=LIGHT, border=BORDER, radius=0.05)
tfm = tb(s, px + Inches(0.35), y + Inches(0.3), SLIDE_W - MR - px - Inches(0.7), Inches(2.1), anchor=MSO_ANCHOR.MIDDLE)
p = para(tfm, 15, MUTED, bold=True, first=True, after=8)
run(p, "FULL PHYSICS CATALOGUE", 15, MUTED, bold=True)
p = para(tfm, 22, ORANGE, bold=True, after=2, leading=1.0)
run(p, "Just us:  ~years", 22, ORANGE, bold=True)
p = para(tfm, 22, BLUE, bold=True, after=0, leading=1.0)
run(p, "With a team:  ~2 yrs", 22, BLUE, bold=True)
pagenum(s, 8)
notes(s, "The funding logic, shown honestly. Two people can PROVE it but cannot build the catalogue the market needs in "
         "any reasonable time. So the plan is: prove with 2 → then a small, capital-efficient India team multiplies "
         "output (~₹50–80L/yr in year 2). This is why a raise makes sense AFTER the pilot — the money buys content "
         "velocity, the one real bottleneck.")

# ============================================================================= SLIDE 9 — 3-YEAR ROADMAP
s = slide(); bg(s, WHITE)
y = header(s, "08 · The 3-year journey", "India → UAE → US. Prove first, then scale.")
gap = Inches(0.34); cw3 = (CW - 2*gap) / 3; ch = Inches(3.7)
def year_card(x, tag, tag_color, heading, items):
    card(s, x, y, cw3, ch, fill=LIGHT, border=BORDER, radius=0.05)
    rect(s, x, y, cw3, Inches(0.12), tag_color, radius=None)
    tf = tb(s, x + Inches(0.3), y + Inches(0.34), cw3 - Inches(0.6), ch - Inches(0.6))
    p = para(tf, 14, tag_color, bold=True, first=True, after=2)
    run(p, tag, 14, tag_color, bold=True)
    p = para(tf, 17, INK, bold=True, after=10, leading=1.02)
    run(p, heading, 17, INK, bold=True)
    for it in items:
        b = para(tf, 12.5, INK, after=7, leading=1.1)
        run(b, "•  ", 12.5, tag_color, bold=True); run(b, it, 12.5, INK)
year_card(ML, "YEAR 1 — PROVE", GREEN, "India · prove it",
          ["Founder + IIT professor", "Build the physics pilot chapter",
           "1–2 institutes, free", "Prove the lift on 100–150 students",
           "End of Y1: start selling to institutes"])
year_card(ML + cw3 + gap, "YEAR 2 — SCALE", BLUE, "India + UAE · grow",
          ["Small content team begins", "Finish physics, start maths + chemistry",
           "Launch UAE — same syllabus, zero rebuild", "Add the AI chatbot → voice tutor",
           "First real revenue at scale"])
year_card(ML + 2*(cw3 + gap), "YEAR 3 — EXPAND", ORANGE, "Go global",
          ["Scale the catalogue with the team", "Enter US / UK (after syllabus rebuild)",
           "Personalised AI tutor matures", "Quality compounds week by week"])
tf2 = tb(s, ML, y + ch + Inches(0.12), CW, Inches(0.4))
p = para(tf2, 13, MUTED, first=True, after=0, align=PP_ALIGN.CENTER)
run(p, "A team makes every step faster.", 13, MUTED, italic=True)
pagenum(s, 9)
notes(s, "The journey. Year 1 is cheap and about PROOF — two of us, one chapter, 1–2 institutes free, measure the "
         "comprehension lift. End of year 1 we start selling. Year 2: a small team, finish physics + start maths/chem, "
         "and open the UAE (same CBSE, zero rebuild) with the AI chatbot/voice tutor. Year 3: scale and enter the West "
         "after a syllabus rebuild. Each stage de-risks the next; a team accelerates all of it.")

# ============================================================================= SLIDE 10 — THE NEXT STEP
s = slide(); bg(s, BG_DARK)
rect(s, 0, 0, SLIDE_W, Inches(0.16), AMBER)
rect(s, 0, 0, Inches(0.22), SLIDE_H, AMBER)
tf = tb(s, Inches(1.1), Inches(0.7), Inches(11), Inches(1.4))
p = para(tf, 14, AMBER, bold=True, first=True, after=4)
run(p, "09 · WHAT'S NEXT", 14, AMBER, bold=True)
p = para(tf, 38, WHITE, bold=True, leading=1.0, after=0)
run(p, "The next step: the pilot.", 38, WHITE, bold=True)
tf2 = tb(s, Inches(1.1), Inches(2.45), Inches(11), Inches(3.4))
bullet_rich(tf2, [("The product is live. The method works. ", WHITE, True),
                  ("The one thing left is proof.", ONDARK, False)],
            lead_color=AMBER, size=18, first=True, after=16)
bullet_rich(tf2, [("Next: run the physics pilot with an institute — ", WHITE, True),
                  ("100–150 students, measuring how much more they understand, before vs after.", ONDARK, False)],
            lead_color=AMBER, size=18, after=16)
bullet_rich(tf2, [("That proof unlocks everything on the roadmap — ", WHITE, True),
                  ("distribution, the team, and the journey abroad.", ONDARK, False)],
            lead_color=AMBER, size=18, after=0)
rect(s, Inches(1.12), Inches(6.15), Inches(3.2), Pt(2), AMBER)
tf3 = tb(s, Inches(1.12), Inches(6.35), Inches(11), Inches(0.6))
p = para(tf3, 14, RGBColor(0xCB,0xD5,0xE1), first=True, after=0)
run(p, "Pradeep  ·  ", 14, WHITE, bold=True)
run(p, "nagapuripradeep02@gmail.com", 14, ONDARK)
notes(s, "Close as forward motion, NOT a request. Don't ask for money or for access. Say: the product is real, the "
         "method works, and the natural next step is a small pilot to PROVE the learning lift — and that proof unlocks "
         "everything after it. Let the listener volunteer to help. Then stop and take questions.")

# ----------------------------------------------------------------------------- save
out = r"C:/Tutor/physics-mind/docs/pitch/PhysicsMind_Pitch.pptx"
prs.save(out)
print("Saved:", out, "—", len(prs.slides._sldIdLst), "slides")
